import fitz
import textract
import docx
from pptx import Presentation


class DocumentProcessor:
    def __init__(self):
        pass

    def extract_text(self, uploaded_file):
        suffix = uploaded_file.name.split(".")[-1]
        if suffix == "pdf": full_text = self.extract_pdf_text(uploaded_file)
        elif suffix == "docx" or suffix == "doc": full_text = self.extract_word_text(uploaded_file)
        elif suffix == "txt": full_text = self.extract_txt_text(uploaded_file)
        elif suffix == "pptx": full_text = self.extract_pptx_text(uploaded_file)
        else: raise ValueError("Unsupported file type. Please upload a PDF file.")

        return full_text
    
    def extract_word_text(self, uploaded_file):
        doc = docx.Document(uploaded_file)
        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text)
        return '\n'.join(full_text)

    def extract_txt_text(self, uploaded_file):
        text = uploaded_file.getvalue().decode("utf-8").replace("\n", " ").replace("  ", " ")
        return text

    def extract_pptx_text(self, uploaded_file):
        ppt = Presentation(uploaded_file)
        full_text = ""
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text += shape.text + " "
        return full_text
    

    def extract_pdf_text(self, uploaded_file):
        pdf_document = fitz.open(stream=uploaded_file.read(), filetype="pdf")
        full_text = ""
        for page_number in range(pdf_document.page_count):
            page = pdf_document.load_page(page_number)
            page_text = page.get_text("text")
            full_text += page_text + "\n\n"
        pdf_document.close()
        return full_text